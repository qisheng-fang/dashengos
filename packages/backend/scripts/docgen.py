#!/usr/bin/env python3
"""
DaShengOS Document Generator — Phase A.4 (v0.3)
Reads JSON from stdin, generates PPTX/DOCX/PDF/XLSX, prints output path to stdout.

Usage:
  echo '{"format":"pptx","title":"My Deck","slides":[{"title":"S1","content":"..."}]}' | python3 docgen.py

Output (stdout):  {"ok":true,"path":"/tmp/dasheng-docs/doc_abc123.pptx","size":12345}
Error  (stderr):  "ERROR: ..."  (exit code 1)

Supported formats:
  - pptx: slides with title/content (python-pptx)
  - docx: sections with heading/content (python-docx)
  - pdf:  HTML content → PDF via weasyprint
  - xlsx: sheets with headers/rows (openpyxl)
"""

import json
import sys
import os
import time
import uuid

OUTPUT_DIR = "/tmp/dasheng-docs"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ─── Helpers ────────────────────────────────────────────────

def make_path(fmt: str) -> str:
    """Generate a unique output file path."""
    uid = uuid.uuid4().hex[:10]
    ext = {"pptx": ".pptx", "docx": ".docx", "pdf": ".pdf", "xlsx": ".xlsx"}[fmt]
    return os.path.join(OUTPUT_DIR, f"doc_{uid}{ext}")

def clean_html(html: str) -> str:
    """Wrap bare HTML into a minimal document for weasyprint."""
    if html.strip().startswith("<!DOCTYPE") or html.strip().startswith("<html"):
        return html
    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head><meta charset="utf-8"><style>
  body {{ font-family: -apple-system, "Microsoft YaHei", sans-serif; max-width: 800px; margin: 40px auto; padding: 20px; color: #1a1a1a; line-height: 1.8; }}
  h1 {{ font-size: 28px; border-bottom: 2px solid #FF6B35; padding-bottom: 8px; }}
  h2 {{ font-size: 22px; margin-top: 28px; color: #FF6B35; }}
  h3 {{ font-size: 18px; }}
  p {{ margin: 12px 0; }}
  ul, ol {{ margin: 8px 0 8px 24px; }}
  li {{ margin: 4px 0; }}
  table {{ border-collapse: collapse; width: 100%; margin: 16px 0; }}
  th, td {{ border: 1px solid #ddd; padding: 8px 12px; text-align: left; }}
  th {{ background: #f5f5f5; font-weight: 600; }}
  code {{ background: #f0f0f0; padding: 2px 6px; border-radius: 3px; font-size: 0.9em; }}
  pre {{ background: #f5f5f5; padding: 16px; border-radius: 6px; overflow-x: auto; }}
  .page-break {{ page-break-after: always; }}
</style></head>
<body>
{html}
</body></html>"""

# ─── PPTX Generator ─────────────────────────────────────────

def generate_pptx(data: dict, path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = data.get("title", "Untitled")
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = "DaShengOS AI 生成"

    # Content slides
    for item in data.get("slides", []):
        # Use blank layout for full control
        blank_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(blank_layout)

        # Title box
        left = Inches(0.8)
        top = Inches(0.4)
        width = Inches(11.7)
        height = Inches(1.0)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0x6B, 0x35)

        # Content box
        left = Inches(0.8)
        top = Inches(1.6)
        width = Inches(11.7)
        height = Inches(5.2)
        content_box = slide.shapes.add_textbox(left, top, width, height)
        tf = content_box.text_frame
        tf.word_wrap = True

        content = item.get("content", "")
        lines = content.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            if line.startswith("# "):
                p.font.size = Pt(22)
                p.font.bold = True
            elif line.startswith("## "):
                p.font.size = Pt(18)
                p.font.bold = True
            elif line.startswith("- ") or line.startswith("* "):
                p.font.size = Pt(16)
                p.level = 0
            else:
                p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_after = Pt(6)

    prs.save(path)


# ─── DOCX Generator ─────────────────────────────────────────

def generate_docx(data: dict, path: str) -> None:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Title
    title_para = doc.add_heading(data.get("title", "Untitled"), level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph("")  # spacer

    # Sections
    for sec in data.get("sections", []):
        doc.add_heading(sec.get("heading", ""), level=1)

        content = sec.get("content", "")
        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                doc.add_paragraph("")
            elif stripped.startswith("- ") or stripped.startswith("* "):
                doc.add_paragraph(stripped[2:], style="List Bullet")
            elif stripped.startswith("1. ") or stripped.startswith("2. "):
                doc.add_paragraph(stripped[3:], style="List Number")
            elif stripped.startswith("```"):
                continue  # skip code fence markers
            else:
                doc.add_paragraph(stripped)

        doc.add_paragraph("")  # spacer between sections

    doc.save(path)


# ─── PDF Generator (weasyprint) ─────────────────────────────

def generate_pdf(data: dict, path: str) -> None:
    try:
        from weasyprint import HTML
    except ImportError:
        # Graceful fallback message
        msg = (
            "weasyprint 未安装。请运行:\n"
            "  pip install weasyprint\n"
            "或安装系统依赖:\n"
            "  brew install pango cairo (macOS)\n"
            "  apt install libpango-1.0-0 libcairo2 (Linux)"
        )
        print(json.dumps({"ok": False, "error": msg}))
        sys.exit(1)

    html = clean_html(data.get("html", ""))
    doc = HTML(string=html)
    doc.write_pdf(path)


# ─── XLSX Generator ─────────────────────────────────────────

def generate_xlsx(data: dict, path: str) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = Workbook()

    # Remove default sheet
    wb.remove(wb.active)

    header_font = Font(bold=True, color="FFFFFF", size=12)
    header_fill = PatternFill(start_color="FF6B35", end_color="FF6B35", fill_type="solid")
    header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    cell_align = Alignment(vertical="top", wrap_text=True)
    thin_border = Border(
        left=Side(style="thin", color="CCCCCC"),
        right=Side(style="thin", color="CCCCCC"),
        top=Side(style="thin", color="CCCCCC"),
        bottom=Side(style="thin", color="CCCCCC"),
    )

    for sheet_data in data.get("sheets", []):
        ws_name = sheet_data.get("name", "Sheet")[:31]  # Excel sheet name limit
        ws = wb.create_sheet(title=ws_name)

        headers = sheet_data.get("headers", [])
        rows = sheet_data.get("rows", [])

        # Write headers
        for ci, h in enumerate(headers, start=1):
            cell = ws.cell(row=1, column=ci, value=h)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_align
            cell.border = thin_border

        # Write data rows
        for ri, row in enumerate(rows, start=2):
            for ci, val in enumerate(row, start=1):
                cell = ws.cell(row=ri, column=ci, value=val)
                cell.alignment = cell_align
                cell.border = thin_border

        # Auto-width columns (approximate)
        for col_cells in ws.columns:
            max_len = 0
            col_letter = col_cells[0].column_letter
            for cell in col_cells:
                val_len = len(str(cell.value or ""))
                # CJK characters are roughly 2x width
                cjk_count = sum(1 for c in str(cell.value or "") if '\u4e00' <= c <= '\u9fff')
                effective_len = val_len + cjk_count  # double-count CJK
                if effective_len > max_len:
                    max_len = effective_len
            ws.column_dimensions[col_letter].width = min(max_len + 4, 60)

    wb.save(path)


# ─── Auto-cleanup (remove files older than 1 hour) ──────────

def cleanup_old_files():
    """Remove files in OUTPUT_DIR older than 1 hour."""
    now = time.time()
    max_age = 3600  # 1 hour
    try:
        for fname in os.listdir(OUTPUT_DIR):
            fpath = os.path.join(OUTPUT_DIR, fname)
            if os.path.isfile(fpath):
                age = now - os.path.getmtime(fpath)
                if age > max_age:
                    os.remove(fpath)
    except OSError:
        pass  # best-effort


# ─── Main ───────────────────────────────────────────────────

def main():
    cleanup_old_files()

    try:
        raw = sys.stdin.read()
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid JSON input: {e}", file=sys.stderr)
        sys.exit(1)

    fmt = data.get("format", "").lower()
    if fmt not in ("pptx", "docx", "pdf", "xlsx"):
        print(f"ERROR: Unsupported format '{fmt}'. Supported: pptx, docx, pdf, xlsx", file=sys.stderr)
        sys.exit(1)

    path = make_path(fmt)

    try:
        if fmt == "pptx":
            generate_pptx(data, path)
        elif fmt == "docx":
            generate_docx(data, path)
        elif fmt == "pdf":
            generate_pdf(data, path)
        elif fmt == "xlsx":
            generate_xlsx(data, path)
    except ImportError as e:
        lib_map = {"pptx": "python-pptx", "docx": "python-docx", "pdf": "weasyprint", "xlsx": "openpyxl"}
        pkg = lib_map.get(fmt, fmt)
        err = {
            "error": f"{pkg} 未安装。请运行: pip install {pkg}",
            "missing_package": pkg,
        }
        print(json.dumps(err))
        sys.exit(1)
    except Exception as e:
        print(f"ERROR: {type(e).__name__}: {e}", file=sys.stderr)
        sys.exit(1)

    size = os.path.getsize(path)
    result = {"ok": True, "path": path, "file_name": os.path.basename(path), "size": size}
    print(json.dumps(result))


if __name__ == "__main__":
    main()
